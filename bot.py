# bot.py | branch: styled
import os, asyncio, tempfile, re
from aiogram import Bot, Dispatcher, F
from aiogram.types import Message, FSInputFile
from aiogram.filters import CommandStart
from openai import AsyncOpenAI
from pydub import AudioSegment
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

load_dotenv()

bot = Bot(token=os.getenv("TELEGRAM_TOKEN"))
dp = Dispatcher()
client = AsyncOpenAI(
    api_key=os.getenv("OPENAI_API_KEY"),
    base_url="https://api.proxyapi.ru/openai/v1"
)

# ── Палитра ────────────────────────────────────────────────────────
BG_BLACK  = RGBColor(0x00, 0x00, 0x00)  # титул, итоги, перебивки
BG_LIGHT  = RGBColor(0xEC, 0xF1, 0xF3)  # контентные слайды
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_BLACK = RGBColor(0x00, 0x00, 0x00)
CLR_GRAY  = RGBColor(0x7C, 0x7C, 0x7C)
CLR_ORANGE = RGBColor(0xFF, 0x5C, 0x00)
CLR_PURPLE = RGBColor(0x7B, 0x4C, 0xFF)
CLR_TEAL   = RGBColor(0x02, 0xB2, 0xA9)
CLR_BLUE   = RGBColor(0x45, 0x79, 0xFF)
ACCENT_CYCLE = [CLR_PURPLE, CLR_ORANGE, CLR_TEAL, CLR_BLUE]
FONT = "Helvetica Neue"

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.625)
LABEL   = "Москва 2026"

# ── Состояние пользователей ────────────────────────────────────────
# format: "voice" | "audio" | "text"
# timing: int (5–120)
# raw_input: исходный текст (из речи/файла/описания)
# md_plan: финальный Markdown для PPTX
state = {}

# ── Системный промт с таймингом и развёрнутыми буллетами ───────────

SYSTEM_PROMPT_BASE = """
Ты — ассистент для создания презентаций.
На вход ты получаешь сырой контент (расшифровка речи или текстовое описание).
Твоя задача —:
1) Сначала спланировать структуру презентации под заданный тайминг.
2) Затем выдать финальную структуру в формате Markdown-презентации.

Правила структуры:
- Первый слайд — титульный: # Заголовок + одна строка подзаголовка буллетом.
- Контентные слайды: ## Заголовок раздела + 3–5 буллетов.
- Последний слайд — итоги: ### Ключевые выводы + 3–5 буллетов.
- Каждый буллет — одно полноценное предложение 10–20 слов,
  которое раскрывает мысль: без телеграфного стиля и обрывков фраз.
- Старайся распределять материал равномерно по слайдам.
- Используй ТОЛЬКО то, что сказал пользователь, ничего не придумывай от себя.

Тайминг:
{timing_hint}

Формат вывода (строго):
ПЕРВОЙ частью выведи раздел:

План:
1. Слайд 1 — ...
2. Слайд 2 — ...
...

Затем пустая строка и раздел:

Презентация:
# Заголовок титульного
- Подзаголовок как одно предложение

## Заголовок раздела
- Развёрнутый буллет номер один
- Развёрнутый буллет номер два

### Ключевые выводы
- ...

Никаких комментариев до блока «План:» и после блока «Презентация:».
"""

def timing_hint_from_minutes(minutes: int) -> str:
    if minutes <= 20:
        return f"- {minutes} минут: 8–10 слайдов, ~2 минуты на слайд, буллеты максимально ёмкие."
    if minutes <= 35:
        return f"- {minutes} минут: 10–14 слайдов, ~2–3 минуты на слайд, можно чуть больше деталей."
    return f"- {minutes} минут: 14–18+ слайдов, больше примеров и расшифровок, но без воды."


# ── Утилиты стилизации ─────────────────────────────────────────────

def set_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=CLR_WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color

def add_rect(slide, left, top, width, height, color: RGBColor):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def add_chrome(slide, slide_num: int, dark: bool):
    bar_color = RGBColor(0x1A, 0x1A, 0x1A) if dark else RGBColor(0xD8, 0xDE, 0xE1)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.38), bar_color)
    label_color = CLR_WHITE if dark else CLR_GRAY
    add_text(slide, "avito.tech", Inches(0.14), Inches(0.05), Inches(1.2), Inches(0.32),
             size=10, bold=True, color=label_color)
    add_text(slide, LABEL, Inches(1.4), Inches(0.05), Inches(2.5), Inches(0.32),
             size=10, bold=False, color=label_color)
    num_color = CLR_WHITE if dark else CLR_GRAY
    add_text(slide, str(slide_num),
             Inches(9.3), Inches(0.03), Inches(0.5), Inches(0.34),
             size=13, bold=True, color=num_color, align=PP_ALIGN.RIGHT)


# ── Три типа слайдов ───────────────────────────────────────────────

def make_dark_slide(prs, title: str, bullets: list, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_BLACK)
    add_chrome(slide, slide_num, dark=True)
    add_rect(slide, Inches(0), Inches(0.38), Inches(0.06), Inches(5.245), CLR_ORANGE)
    add_text(slide, title,
             Inches(0.2), Inches(0.65), Inches(8.5), Inches(1.6),
             size=48, bold=True, color=CLR_WHITE)
    for i, bullet in enumerate(bullets[:5]):
        top = Inches(2.4) + i * Inches(0.6)
        add_text(slide, bullet,
                 Inches(0.2), top, Inches(8.5), Inches(0.55),
                 size=18, bold=False, color=CLR_WHITE)

def make_content_slide(prs, title: str, bullets: list, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    add_chrome(slide, slide_num, dark=False)
    add_text(slide, title,
             Inches(0.22), Inches(0.5), Inches(9.3), Inches(0.85),
             size=32, bold=True, color=CLR_BLACK)
    add_rect(slide, Inches(0.22), Inches(1.38), Inches(9.1), Inches(0.04), CLR_ORANGE)
    for i, bullet in enumerate(bullets[:5]):
        top = Inches(1.58) + i * Inches(0.78)
        accent = ACCENT_CYCLE[i % 4]
        add_rect(slide, Inches(0.22), top + Inches(0.1), Inches(0.07), Inches(0.26), accent)
        add_text(slide, bullet,
                 Inches(0.38), top, Inches(9.1), Inches(0.68),
                 size=16, bold=False, color=CLR_BLACK)


# ── Markdown → PPTX ────────────────────────────────────────────────

def build_pptx(md_content: str, output_path: str):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = []
    current_title   = None
    current_bullets = []
    current_level   = 1

    def flush():
        nonlocal current_title, current_bullets, current_level
        if current_title is None:
            return
        slides.append({
            "title":   current_title,
            "bullets": list(current_bullets),
            "level":   current_level,
        })

    for line in md_content.splitlines():
        line = line.strip()
        if line.startswith("# ") and not line.startswith("## "):
            flush()
            current_title   = line[2:].strip()
            current_bullets = []
            current_level   = 1
        elif line.startswith("## ") and not line.startswith("### "):
            flush()
            current_title   = line[3:].strip()
            current_bullets = []
            current_level   = 2
        elif line.startswith("### "):
            flush()
            current_title   = line[4:].strip()
            current_bullets = []
            current_level   = 3
        elif line.startswith(("- ", "* ", "• ")):
            current_bullets.append(line[2:].strip())

    flush()

    for i, s in enumerate(slides):
        num = i + 1
        is_dark = (s["level"] in (1, 3))
        if is_dark:
            make_dark_slide(prs, s["title"], s["bullets"], num)
        else:
            make_content_slide(prs, s["title"], s["bullets"], num)

    prs.save(output_path)


# ── GPT: план + Markdown по таймингу ───────────────────────────────

async def build_plan_and_markdown(user_id: int):
    st = state[user_id]
    minutes = st["timing"]
    system_prompt = SYSTEM_PROMPT_BASE.format(
        timing_hint=timing_hint_from_minutes(minutes)
    )
    resp = await client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": st["raw_input"]}
        ]
    )
    full = resp.choices[0].message.content.strip()

    plan_part = ""
    md_part   = full
    if "Презентация:" in full:
        parts = full.split("Презентация:", 1)
        plan_part = parts[0].strip()
        md_part   = parts[1].strip()

    st["md_plan"] = md_part

    slides_plan = []
    for line in plan_part.splitlines():
        line = line.strip()
        m = re.match(r"^(\d+)[\.\)]\s*(.+)$", line)
        if not m:
            continue
        idx = int(m.group(1))
        title = m.group(2)
        slides_plan.append((idx, title))

    return plan_part, md_part, slides_plan


# ── Telegram-хэндлеры ─────────────────────────────────────────────-

@dp.message(CommandStart())
async def start(message: Message):
    user_id = message.from_user.id
    state[user_id] = {
        "format": None,
        "timing": None,
        "raw_input": "",
        "md_plan": ""
    }
    await message.answer(
        "Привет! 🎤\n\n"
        "Я помогу собрать презентацию.\n\n"
        "Выбери формат ввода:\n"
        "1 — голосовые сообщения\n"
        "2 — загрузить аудиофайл\n"
        "3 — текстовое описание\n\n"
        "Просто ответь: 1, 2 или 3."
    )


@dp.message(F.text.in_(["1", "2", "3"]))
async def choose_format(message: Message):
    user_id = message.from_user.id
    st = state.setdefault(user_id, {"format": None, "timing": None, "raw_input": "", "md_plan": ""})

    if message.text == "1":
        st["format"] = "voice"
        await message.answer(
            "Формат: голосовые сообщения.\n"
            "Теперь отправь длительность выступления в минутах (например, 20 или 35)."
        )
    elif message.text == "2":
        st["format"] = "audio"
        await message.answer(
            "Формат: аудиофайл.\n"
            "Сначала отправь длительность выступления в минутах (например, 25).\n"
            "После этого пришли файл (.mp3, .m4a, .wav, .ogg)."
        )
    else:
        st["format"] = "text"
        await message.answer(
            "Формат: текстовое описание.\n"
            "Отправь длительность выступления в минутах (например, 30).\n"
            "После выбора просто пришли текст/описание."
        )


@dp.message(F.text)
async def handle_timing_or_commands(message: Message):
    user_id = message.from_user.id
    text = message.text.strip()

    # Служебные команды обрабатываются отдельными хэндлерами ниже
    if text in ["/plan", "/build", "/start", "1", "2", "3"]:
        return

    # Пытаемся распарсить как число (тайминг)
    if text.isdigit():
        minutes = int(text)
        if minutes < 5 or minutes > 120:
            await message.answer("Введи длительность от 5 до 120 минут (например, 25).")
            return

        st = state.setdefault(user_id, {"format": None, "timing": None, "raw_input": "", "md_plan": ""})
        st["timing"] = minutes

        if not st["format"]:
            await message.answer(
                f"Тайминг {minutes} минут сохранён.\n"
                "Теперь выбери формат: 1 — голос, 2 — аудиофайл, 3 — текст."
            )
            return

        if st["format"] == "voice":
            await message.answer(
                f"Ок, делаем презентацию на {minutes} минут по голосовым.\n\n"
                "Наговаривай одно или несколько голосовых.\n"
                "Когда закончишь — напиши /plan."
            )
        elif st["format"] == "audio":
            await message.answer(
                f"Ок, делаем презентацию на {minutes} минут по аудиофайлу.\n\n"
                "Пришли аудиофайл (mp3/m4a/wav/ogg), потом напиши /plan."
            )
        else:
            await message.answer(
                f"Ок, делаем презентацию на {minutes} минут по тексту.\n\n"
                "Пришли текст/описание одним или несколькими сообщениями, затем /plan."
            )
        return

    # Если не число — считаем текстовым вводом (формат "text")
    st = state.setdefault(user_id, {"format": None, "timing": None, "raw_input": "", "md_plan": ""})
    if st.get("format") != "text":
        return
    if not st.get("timing"):
        await message.answer("Сначала задай тайминг: просто отправь количество минут (например, 25).")
        return

    st["raw_input"] = (st.get("raw_input") or "") + "\n\n" + text
    await message.answer(
        "✍️ Текст добавлен.\n"
        "Можешь дописать ещё или отправить /plan для предварительной структуры."
    )


# ── Голосовые сообщения ────────────────────────────────────────────

@dp.message(F.voice)
async def handle_voice(message: Message):
    user_id = message.from_user.id
    st = state.setdefault(user_id, {"format": "voice", "timing": None, "raw_input": "", "md_plan": ""})
    if st.get("format") != "voice":
        await message.answer("Сейчас выбран другой формат. Чтобы перейти на голос, отправь /start и выбери 1.")
        return
    if not st.get("timing"):
        await message.answer("Сначала задай тайминг: просто отправь количество минут (например, 25).")
        return

    await message.answer("⏳ Транскрибирую голос...")
    voice_file = await bot.get_file(message.voice.file_id)

    with tempfile.TemporaryDirectory() as tmp:
        oga_path = f"{tmp}/voice.oga"
        mp3_path = f"{tmp}/voice.mp3"
        await bot.download_file(voice_file.file_path, oga_path)
        AudioSegment.from_ogg(oga_path).export(mp3_path, format="mp3")

        with open(mp3_path, "rb") as f:
            result = await client.audio.transcriptions.create(
                model="whisper-1", file=f, language="ru"
            )

    text = result.text
    st["raw_input"] = (st.get("raw_input") or "") + "\n\n" + text
    await message.answer(
        f"✅ Записал фрагмент:\n\n_{text}_\n\n"
        "Наговори ещё или напиши /plan для предварительной структуры.",
        parse_mode="Markdown"
    )


# ── Аудиофайлы ─────────────────────────────────────────────────────

@dp.message(F.audio | F.document)
async def handle_audio_file(message: Message):
    user_id = message.from_user.id
    st = state.setdefault(user_id, {"format": None, "timing": None, "raw_input": "", "md_plan": ""})
    if st.get("format") != "audio":
        return

    if not st.get("timing"):
        await message.answer("Сначала задай тайминг: просто отправь количество минут (например, 25).")
        return

    telegram_file = message.audio or message.document
    await message.answer("⏳ Транскрибирую аудиофайл...")

    file_info = await bot.get_file(telegram_file.file_id)
    with tempfile.TemporaryDirectory() as tmp:
        in_path = f"{tmp}/input"
        out_path = f"{tmp}/audio.mp3"
        await bot.download_file(file_info.file_path, in_path)
        AudioSegment.from_file(in_path).export(out_path, format="mp3")

        with open(out_path, "rb") as f:
            result = await client.audio.transcriptions.create(
                model="whisper-1", file=f, language="ru"
            )

    text = result.text
    st["raw_input"] = text
    await message.answer(
        "✅ Аудиофайл распознан.\n\n"
        "Теперь напиши /plan, чтобы я показал структуру слайдов."
    )


# ── Предварительный план ──────────────────────────────────────────

@dp.message(F.text == "/plan")
async def show_plan(message: Message):
    user_id = message.from_user.id
    st = state.get(user_id)
    if not st or not st.get("raw_input"):
        await message.answer("Мне пока не из чего делать презентацию. Сначала пришли контент.")
        return
    if not st.get("timing"):
        await message.answer("Сначала задай тайминг: просто отправь количество минут (например, 25).")
        return

    await message.answer("🧩 Собираю структуру и тезисы...")
    plan_text, md, slides_plan = await build_plan_and_markdown(user_id)

    minutes = st["timing"]
    total_slides = len(slides_plan) if slides_plan else "?"
    per_slide = round(minutes / total_slides, 1) if isinstance(total_slides, int) and total_slides > 0 else "≈2–3"

    lines = [
        f"План презентации на {minutes} минут:",
        f"Слайдов: {total_slides}, ориентировочно {per_slide} мин на слайд.",
        ""
    ]
    if slides_plan:
        for idx, title in slides_plan:
            lines.append(f"{idx}. {title}")
    else:
        lines.append(plan_text or "План неявно вшит в Markdown.")

    lines.append("")
    lines.append("Если структура ок — напиши /build, и я соберу PPTX.")
    lines.append("Если что-то не так — дополни входной текст и снова вызови /plan.")

    await message.answer("\n".join(lines))


# ── Сборка презентации ────────────────────────────────────────────

@dp.message(F.text == "/build")
async def build_presentation(message: Message):
    user_id = message.from_user.id
    st = state.get(user_id)
    if not st or not st.get("raw_input"):
        await message.answer("Сначала пришли контент (голос, аудио или текст).")
        return
    if not st.get("timing"):
        await message.answer("Сначала задай тайминг: просто отправь количество минут (например, 25).")
        return

    if not st.get("md_plan"):
        await message.answer("Сначала делаю структуру...")
        _, _, _ = await build_plan_and_markdown(user_id)

    await message.answer("🔨 Собираю презентацию...")

    md_content = st["md_plan"]

    with tempfile.TemporaryDirectory() as tmp:
        pptx_path = f"{tmp}/presentation_{st['timing']}min.pptx"
        build_pptx(md_content, pptx_path)
        await message.answer_document(
            FSInputFile(pptx_path, filename=os.path.basename(pptx_path)),
            caption="🎉 Готово! Презентация на основе твоих материалов."
        )

    state[user_id] = {
        "format": st["format"],
        "timing": st["timing"],
        "raw_input": "",
        "md_plan": ""
    }


async def main():
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
